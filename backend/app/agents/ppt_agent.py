"""PPT Report Agent — generates an 8-slide PowerPoint report from insight_report + raw metrics.

Input  (PipelineState keys consumed):
    insight_report      : dict   — InsightReport.model_dump() from insight_agent
    performance_metrics : dict   — from performance_agent
    funnel_metrics      : dict   — from funnel_agent
    cohort_metrics      : dict   — from cohort_agent
    journey_metrics     : dict   — from journey_agent
    anomaly_metrics     : dict   — from anomaly_agent
    prediction_metrics  : dict   — from prediction_agent
    domain_context      : dict   — DomainContext.model_dump() from context_agent
    week_start          : str    — "YYYYMMDD"
    week_end            : str    — "YYYYMMDD"

Output (PipelineState key produced):
    ppt_url : str   — local file path (or object-storage URL) of the generated .pptx

Slide structure (fixed 8 slides):
    Slide 1  Executive Summary      — KPI 카드 3개 + 한 줄 요약
    Slide 2  주요 지표 현황           — KPI 테이블 + insight chart_type 기반 시각화
    Slide 3  이상 감지 결과           — 이상값 강조 테이블
    Slide 4  사용자 흐름 분석         — 퍼널 단계별 이탈 + insight chart_type 기반 시각화
    Slide 5  고객 세그먼트 분석        — 디바이스/소스별 비교
    Slide 6  [도메인 가변 슬라이드]    — e-commerce: 카테고리별 구매 분석
    Slide 7  예측 및 시사점           — insight chart_type 기반 시각화 + LLM 코멘트
    Slide 8  권장 액션               — LLM recommendations + 우선순위
"""

from __future__ import annotations

import logging
import os
from datetime import datetime
from pathlib import Path
from typing import Any

from lxml import etree
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

logger = logging.getLogger(__name__)

# ──────────────────────────────────────────────────────────────────────────────
# Layout / theme constants  (bright mode)
# ──────────────────────────────────────────────────────────────────────────────

_SLIDE_W = Inches(13.33)
_SLIDE_H = Inches(7.5)

_C = {
    "bg":          RGBColor(0xF7, 0xF9, 0xFC),   # 슬라이드 배경 (연회색)
    "card":        RGBColor(0xFF, 0xFF, 0xFF),   # 카드/섹션 배경 (흰색)
    "header_bg":   RGBColor(0x1B, 0x4F, 0x72),   # 헤더 바 (딥 네이비)
    "accent":      RGBColor(0x21, 0x8A, 0xBB),   # 강조색 (청록)
    "row_even":    RGBColor(0xEA, 0xF2, 0xFB),   # 테이블 짝수 행
    "row_odd":     RGBColor(0xFF, 0xFF, 0xFF),   # 테이블 홀수 행
    "tbl_header":  RGBColor(0x1B, 0x4F, 0x72),   # 테이블 헤더 bg
    "positive":    RGBColor(0x1E, 0x8B, 0x4C),   # 증가 초록
    "negative":    RGBColor(0xC0, 0x39, 0x2B),   # 감소/이상 빨강
    "neutral":     RGBColor(0xD3, 0x7C, 0x00),   # 중립 주황
    "text_dark":   RGBColor(0x1A, 0x1A, 0x2E),   # 본문 텍스트
    "text_mid":    RGBColor(0x55, 0x6B, 0x7D),   # 보조 텍스트
    "text_light":  RGBColor(0xFF, 0xFF, 0xFF),   # 헤더/카드 위 흰 텍스트
    "highlight":   RGBColor(0xFA, 0xD7, 0xA0),   # 이상감지 강조 (연주황)
}

_LOGO_PATH = Path(__file__).parents[3] / "sk_ax_logo.webp"
_OUTPUT_DIR = os.environ.get("PPT_OUTPUT_DIR", "/tmp/ppt_reports")

# ──────────────────────────────────────────────────────────────────────────────
# Low-level drawing helpers
# ──────────────────────────────────────────────────────────────────────────────

def _add_rect(slide, left, top, width, height, fill_color: RGBColor):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def _add_textbox(
    slide, text: str, left, top, width, height,
    font_size: int = 12, bold: bool = False,
    color: RGBColor = _C["text_dark"],
    align: PP_ALIGN = PP_ALIGN.LEFT,
    wrap: bool = True,
    v_anchor: str = "top",   # "top" | "middle" | "bottom"
) -> None:
    from pptx.enum.text import MSO_ANCHOR
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    _anchor_map = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}
    tf.vertical_anchor = _anchor_map.get(v_anchor, MSO_ANCHOR.TOP)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color


def _slide_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _C["bg"]


def _add_slide_header(slide, title: str, subtitle: str = "") -> None:
    _add_rect(slide, Inches(0), Inches(0), _SLIDE_W, Inches(1.15), _C["header_bg"])
    _add_textbox(slide, title,
                 Inches(0.3), Inches(0.08), Inches(11), Inches(0.65),
                 font_size=26, bold=True, color=_C["text_light"])
    if subtitle:
        _add_textbox(slide, subtitle,
                     Inches(0.3), Inches(0.72), Inches(11), Inches(0.38),
                     font_size=12, color=RGBColor(0xAA, 0xC4, 0xD8))


def _add_logo(slide) -> None:
    """SK ax 로고 — 슬라이드 왼쪽 하단. 항상 맨 마지막에 호출해야 다른 요소에 가리지 않음."""
    if not _LOGO_PATH.exists():
        return
    try:
        logo_h = Inches(0.52)
        logo_w = Inches(0.93)   # 700:394 비율 유지
        slide.shapes.add_picture(
            str(_LOGO_PATH),
            Inches(0.2),
            _SLIDE_H - logo_h - Inches(0.1),
            width=logo_w,
        )
    except Exception:
        pass


def _add_kpi_card(slide, label: str, value: str, delta: str,
                  left, top, width=Inches(3.5), height=Inches(1.8)) -> None:
    _add_rect(slide, left, top, width, height, _C["card"])
    # left border accent
    _add_rect(slide, left, top, Inches(0.06), height, _C["accent"])
    _add_textbox(slide, label,
                 left + Inches(0.15), top + Inches(0.1), width - Inches(0.25), Inches(0.38),
                 font_size=11, color=_C["text_mid"], v_anchor="middle")
    _add_textbox(slide, value,
                 left + Inches(0.15), top + Inches(0.42), width - Inches(0.25), Inches(0.72),
                 font_size=24, bold=True, color=_C["text_dark"], v_anchor="middle")
    if delta:
        delta_color = (_C["positive"] if delta.startswith("+")
                       else _C["negative"] if delta.startswith("-")
                       else _C["text_mid"])
        _add_textbox(slide, delta,
                     left + Inches(0.15), top + Inches(1.28), width - Inches(0.25), Inches(0.38),
                     font_size=12, color=delta_color)


def _set_cell_anchor(cell, anchor: str = "ctr") -> None:
    """anchor: 'ctr' | 't' | 'b'"""
    cell._tc.get_or_add_tcPr().set("anchor", anchor)


# Per-row height cap so tables don't balloon when there are few rows
_ROW_H_MIN = Inches(0.32)
_ROW_H_MAX = Inches(0.55)
_HDR_H     = Inches(0.40)


def _add_table(
    slide, headers: list[str], rows: list[list[str]],
    left, top, width, height,
    highlight_rows: set[int] | None = None,
    font_size: int = 11,
) -> None:
    if not rows:
        _add_textbox(slide, "(데이터 없음)", left, top, width, height,
                     font_size=11, color=_C["text_mid"])
        return

    n_data = len(rows)
    cols   = len(headers)

    # Dynamic row height: fit within `height`, but clamp per-row
    available = height - _HDR_H
    row_h = max(_ROW_H_MIN, min(_ROW_H_MAX, available // n_data))
    actual_h = _HDR_H + row_h * n_data

    table = slide.shapes.add_table(n_data + 1, cols, left, top, width, actual_h).table
    col_w = width // cols
    for i in range(cols):
        table.columns[i].width = col_w

    # Fix row heights
    table.rows[0].height = _HDR_H
    for r in range(n_data):
        table.rows[r + 1].height = row_h

    # Header row — middle align, bold
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = _C["tbl_header"]
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0] if p.runs else p.add_run()
        run.font.bold = True
        run.font.color.rgb = _C["text_light"]
        run.font.size = Pt(font_size)
        _set_cell_anchor(cell, "ctr")

    # Data rows — top align so text doesn't float in tall cells
    for r, row in enumerate(rows):
        is_hl = highlight_rows and r in highlight_rows
        bg = _C["highlight"] if is_hl else (_C["row_even"] if r % 2 == 0 else _C["row_odd"])
        txt_color = _C["negative"] if is_hl else _C["text_dark"]
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0] if p.runs else p.add_run()
            run.font.color.rgb = txt_color
            run.font.size = Pt(font_size)
            _set_cell_anchor(cell, "ctr")


def _section_label(slide, text: str, left, top, width=Inches(6.0)) -> None:
    _add_textbox(slide, text, left, top, width, Inches(0.35),
                 font_size=11, bold=True, color=_C["text_mid"])


# ──────────────────────────────────────────────────────────────────────────────
# Native chart helpers (bright theme)
# ──────────────────────────────────────────────────────────────────────────────

def _apply_chart_theme(chart) -> None:
    """White background, dark text for all native charts."""
    chart.font.color.rgb = _C["text_dark"]
    chart.font.size = Pt(10)

    cs = chart.element

    def _white_spPr():
        return etree.fromstring(
            '<c:spPr xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart"'
            ' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            '<a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill>'
            '<a:ln><a:solidFill><a:srgbClr val="D0D8E0"/></a:solidFill></a:ln>'
            '</c:spPr>'
        )

    cs.append(_white_spPr())
    pa = cs.find(".//" + qn("c:plotArea"))
    if pa is not None:
        pa.append(_white_spPr())

    # axis label color → dark
    NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    for ax_tag in (qn("c:valAx"), qn("c:catAx")):
        for ax in cs.findall(".//" + ax_tag):
            txPr = ax.find(qn("c:txPr"))
            if txPr is None:
                txPr = etree.SubElement(ax, qn("c:txPr"))
            for tag in (f"{{{NS_A}}}bodyPr", f"{{{NS_A}}}lstStyle"):
                if txPr.find(tag) is None:
                    etree.SubElement(txPr, tag)
            p_el = txPr.find(f"{{{NS_A}}}p")
            if p_el is None:
                p_el = etree.SubElement(txPr, f"{{{NS_A}}}p")
            pPr = p_el.find(f"{{{NS_A}}}pPr")
            if pPr is None:
                pPr = etree.SubElement(p_el, f"{{{NS_A}}}pPr")
            defRPr = pPr.find(f"{{{NS_A}}}defRPr")
            if defRPr is None:
                defRPr = etree.SubElement(pPr, f"{{{NS_A}}}defRPr")
            sf = etree.SubElement(defRPr, f"{{{NS_A}}}solidFill")
            clr = etree.SubElement(sf, f"{{{NS_A}}}srgbClr")
            clr.set("val", "1A1A2E")


def _add_native_chart(
    slide, xl_type,
    categories: list[str], series_data: dict[str, list[float]],
    left, top, width, height,
) -> None:
    cd = ChartData()
    cd.categories = categories
    for name, vals in series_data.items():
        cd.add_series(name, vals)
    gf = slide.shapes.add_chart(xl_type, left, top, width, height, cd)
    _apply_chart_theme(gf.chart)


# ──────────────────────────────────────────────────────────────────────────────
# Formatting helpers
# ──────────────────────────────────────────────────────────────────────────────

def _fmt(v: Any, prefix: str = "", suffix: str = "", decimals: int = 0) -> str:
    try:
        f = float(v)
        return f"{prefix}{f:,.{decimals}f}{suffix}" if decimals else f"{prefix}{int(f):,}{suffix}"
    except (TypeError, ValueError):
        return str(v)


def _pct(v: Any, d: int = 1) -> str:
    try:
        return f"{float(v) * 100:.{d}f}%"
    except (TypeError, ValueError):
        return str(v)


def _wow(v: Any) -> str:
    if v is None:
        return "—"
    try:
        p = float(v) * 100
        return f"+{p:.1f}%" if p >= 0 else f"{p:.1f}%"
    except (TypeError, ValueError):
        return str(v)


# ──────────────────────────────────────────────────────────────────────────────
# chart_type → visualization dispatcher
# Called by slide builders that delegate to insight_report.chart_type
# ──────────────────────────────────────────────────────────────────────────────

def _render_chart_for_slide(
    slide,
    chart_type: str,
    metrics: dict,
    left, top, width, height,
) -> None:
    """
    Dispatch visualization based on insight_agent's chart_type decision.

    Supported:
        line_chart   → daily revenue/session line (performance)
        funnel_chart → horizontal bar per funnel step
        heatmap      → cohort retention color-table (도형 기반)
        bar_chart    → generic column chart
        table        → no extra chart (table already rendered)
        kpi_cards    → no extra chart (cards already rendered)
        sankey       → fallback to journey path table (not supported natively)
    """
    ct = (chart_type or "").lower()

    if ct in ("line_chart", "kpi_cards"):
        # kpi_cards는 카드가 이미 렌더됐지만 daily breakdown이 있으면 라인 차트 추가
        daily = (metrics or {}).get("daily_breakdown", [])[-7:]
        if not daily:
            return
        dates = [d.get("date", "") for d in daily]
        revenues = [float(d.get("revenue", 0) or 0) for d in daily]
        sessions = [float(d.get("session_count", 0) or 0) for d in daily]
        if any(v > 0 for v in revenues):
            _add_native_chart(slide, XL_CHART_TYPE.LINE,
                              dates, {"매출(₩)": revenues, "세션": sessions},
                              left, top, width, height)

    elif ct == "funnel_chart":
        steps = (metrics or {}).get("steps", [])
        if not steps:
            return
        names = [s.get("event_name", "")[:14] for s in steps]
        counts = [float(s.get("user_count", 0) or 0) for s in steps]
        if any(v > 0 for v in counts):
            _add_native_chart(slide, XL_CHART_TYPE.BAR_CLUSTERED,
                              names, {"유저 수": counts},
                              left, top, width, height)

    elif ct == "heatmap":
        # Cohort retention heatmap — 도형 기반 color grid
        _render_cohort_heatmap(slide, metrics, left, top, width, height)

    elif ct == "bar_chart":
        # Generic: try daily breakdown or first numeric series available
        daily = (metrics or {}).get("daily_breakdown", [])[-7:]
        if daily:
            dates = [d.get("date", "") for d in daily]
            vals = [float(d.get("revenue", 0) or 0) for d in daily]
            if any(v > 0 for v in vals):
                _add_native_chart(slide, XL_CHART_TYPE.COLUMN_CLUSTERED,
                                  dates, {"매출(₩)": vals},
                                  left, top, width, height)

    # table / kpi_cards / sankey → no additional native chart
    # (table already rendered; sankey not supported — path table is fallback)


def _render_cohort_heatmap(slide, cohort_metrics: dict, left, top, width, height) -> None:
    """Cohort retention heatmap.

    cohort_metrics["cohorts"] 구조:
        [{ "cohort_week": "2020-W50", "cohort_size": 14,
           "weeks": [{"week_offset": 0, "retention_rate": 1.0}, ...] }]
    """
    cohorts = [c for c in (cohort_metrics or {}).get("cohorts", [])
               if c.get("cohort_week") != "unknown" and c.get("weeks")]
    if not cohorts:
        _add_textbox(slide, "(코호트 데이터 없음)", left, top, width, height,
                     font_size=11, color=_C["text_mid"])
        return

    # 최대 week_offset 수 파악
    max_offset = max(
        (w.get("week_offset", 0) for c in cohorts for w in c.get("weeks", [])),
        default=0
    )
    n_rows = min(len(cohorts), 6)
    n_cols = min(max_offset + 1, 6)   # W0 ~ W(max_offset)

    LABEL_W = Inches(1.3)
    grid_w  = width - LABEL_W
    cell_w  = grid_w  / n_cols
    cell_h  = height  / (n_rows + 1)   # +1 for header row

    GAP = Emu(18000)  # 셀 간격

    def _heat_color(rate: float) -> RGBColor:
        """0.0 → 연파랑, 1.0 → 진파랑"""
        rate = max(0.0, min(1.0, rate))
        r = int(0xD6 + (0x1B - 0xD6) * rate)
        g = int(0xEA + (0x4F - 0xEA) * rate)
        b = int(0xFF + (0x72 - 0xFF) * rate)
        return RGBColor(r, g, b)

    # ── 헤더 행: W0, W1, W2 ... ───────────────────────────────────────────────
    for j in range(n_cols):
        cx = left + LABEL_W + cell_w * j
        _add_rect(slide, cx, top, cell_w - GAP, cell_h - GAP, _C["tbl_header"])
        _add_textbox(slide, f"W{j}", cx, top, cell_w - GAP, cell_h - GAP,
                     font_size=9, bold=True, color=_C["text_light"],
                     align=PP_ALIGN.CENTER, v_anchor="middle")

    # ── 데이터 행 ─────────────────────────────────────────────────────────────
    for i, cohort in enumerate(cohorts[:n_rows]):
        cy_row = top + cell_h * (i + 1)
        label  = cohort.get("cohort_week", f"W{i}")

        # 코호트 라벨 셀
        _add_rect(slide, left, cy_row, LABEL_W - GAP, cell_h - GAP, _C["row_even"])
        _add_textbox(slide, label, left, cy_row, LABEL_W - GAP, cell_h - GAP,
                     font_size=9, bold=True, color=_C["text_dark"],
                     align=PP_ALIGN.CENTER, v_anchor="middle")

        # week_offset → retention_rate 매핑
        offset_map = {w["week_offset"]: w.get("retention_rate", 0)
                      for w in cohort.get("weeks", [])}

        for j in range(n_cols):
            rate  = float(offset_map.get(j, 0) or 0)
            color = _heat_color(rate)
            cx    = left + LABEL_W + cell_w * j
            _add_rect(slide, cx, cy_row, cell_w - GAP, cell_h - GAP, color)
            label_color = _C["text_light"] if rate > 0.5 else _C["text_dark"]
            _add_textbox(slide, f"{rate*100:.0f}%",
                         cx, cy_row, cell_w - GAP, cell_h - GAP,
                         font_size=9, color=label_color,
                         align=PP_ALIGN.CENTER, v_anchor="middle")


# ──────────────────────────────────────────────────────────────────────────────
# Per-slide builders
# ──────────────────────────────────────────────────────────────────────────────

def _build_slide1_executive_summary(slide, insight_report: dict, performance_metrics: dict) -> None:
    _slide_background(slide)
    _add_slide_header(slide, "Executive Summary", "주간 핵심 성과 요약")

    kpis = (performance_metrics or {}).get("kpis", {})
    wow = (performance_metrics or {}).get("wow_change") or {}

    cards = [
        ("총 매출",  _fmt(kpis.get("total_revenue", 0), prefix="₩"),     _wow(wow.get("total_revenue"))),
        ("전환율",   _pct(kpis.get("conversion_rate", 0)),                 _wow(wow.get("conversion_rate"))),
        ("세션 수",  _fmt(kpis.get("session_count", 0)),                   _wow(wow.get("session_count"))),
    ]
    for i, (label, value, delta) in enumerate(cards):
        _add_kpi_card(slide, label, value, delta,
                      Inches(0.4 + i * 4.3), Inches(1.3), Inches(3.9), Inches(1.8))

    summary = (insight_report or {}).get("executive_summary", "")
    _add_rect(slide, Inches(0.4), Inches(3.3), Inches(12.5), Inches(0.5), _C["accent"])
    _add_textbox(slide, "핵심 요약",
                 Inches(0.6), Inches(3.35), Inches(12.0), Inches(0.38),
                 font_size=11, bold=True, color=_C["text_light"])
    _add_rect(slide, Inches(0.4), Inches(3.8), Inches(12.5), Inches(2.5), _C["card"])
    _add_textbox(slide, summary or "(요약 없음)",
                 Inches(0.6), Inches(3.9), Inches(12.1), Inches(2.2),
                 font_size=13, color=_C["text_dark"], wrap=True)

    top_findings = (insight_report or {}).get("top_findings", [])
    if top_findings:
        y = Inches(6.45)
        findings_text = "  •  ".join(top_findings[:3])
        _add_textbox(slide, findings_text,
                     Inches(0.4), y, Inches(12.5), Inches(0.5),
                     font_size=10, color=_C["text_mid"], wrap=False)


def _build_slide2_performance(slide, performance_metrics: dict, insight_report: dict) -> None:
    _slide_background(slide)
    perf_slide = (insight_report or {}).get("performance_slide") or {}
    chart_type = perf_slide.get("chart_type", "line_chart")
    _add_slide_header(slide, "주요 지표 현황", perf_slide.get("headline", ""))

    kpis = (performance_metrics or {}).get("kpis", {})
    wow = (performance_metrics or {}).get("wow_change") or {}

    kpi_rows = [
        ["총 매출",   _fmt(kpis.get("total_revenue", 0), prefix="₩"),    _wow(wow.get("total_revenue"))],
        ["거래 건수", _fmt(kpis.get("transaction_count", 0)),             _wow(wow.get("transaction_count"))],
        ["세션 수",   _fmt(kpis.get("session_count", 0)),                 _wow(wow.get("session_count"))],
        ["전환율",    _pct(kpis.get("conversion_rate", 0)),               _wow(wow.get("conversion_rate"))],
        ["이탈률",    _pct(kpis.get("bounce_rate", 0)),                   "—"],
        ["ARPU",      _fmt(kpis.get("arpu", 0), prefix="₩", decimals=0), "—"],
    ]
    _section_label(slide, "주요 KPI", Inches(0.4), Inches(1.25))
    _add_table(slide, ["지표", "이번 주", "전주 대비"],
               kpi_rows, Inches(0.4), Inches(1.6), Inches(5.2), Inches(3.8))

    # insight_agent가 결정한 chart_type으로 우측 시각화
    _section_label(slide, "일별 추이", Inches(5.9), Inches(1.25), width=Inches(7.0))
    _render_chart_for_slide(
        slide, chart_type, performance_metrics,
        Inches(5.9), Inches(1.6), Inches(7.0), Inches(4.8),
    )

    bullets = perf_slide.get("bullets", [])
    if bullets:
        _add_textbox(slide, "  •  ".join(bullets[:2]),
                     Inches(0.4), Inches(6.5), Inches(12.5), Inches(0.45),
                     font_size=10, color=_C["text_mid"], wrap=False)


def _build_slide3_anomaly(slide, anomaly_metrics: dict, insight_report: dict) -> None:
    _slide_background(slide)
    anomaly_slide = (insight_report or {}).get("anomaly_slide") or {}
    _add_slide_header(slide, "이상 감지 결과", anomaly_slide.get("headline", ""))

    anomalies = (anomaly_metrics or {}).get("anomalies", [])
    summary = (anomaly_metrics or {}).get("summary", {})

    total = summary.get("total_anomalies", 0)
    affected = ", ".join(summary.get("affected_metrics", [])) or "없음"
    worst_date = summary.get("most_abnormal_date") or "—"

    _add_kpi_card(slide, "감지된 이상 수", str(total), "",
                  Inches(0.4), Inches(1.3), Inches(2.8), Inches(1.4))
    _add_kpi_card(slide, "영향 지표", affected[:30], "",
                  Inches(3.5), Inches(1.3), Inches(5.0), Inches(1.4))
    _add_kpi_card(slide, "가장 이상한 날짜", worst_date, "",
                  Inches(8.8), Inches(1.3), Inches(2.8), Inches(1.4))

    if anomalies:
        rows = [
            [
                a.get("metric", ""), a.get("date", ""),
                _fmt(a.get("observed_value", 0), decimals=2),
                _fmt(a.get("expected_mean", 0), decimals=2),
                f"{a.get('z_score', 0):+.2f}",
                a.get("direction", ""),
            ]
            for a in anomalies[:10]
        ]
        highlight = {i for i, a in enumerate(anomalies[:10]) if abs(a.get("z_score", 0)) >= 3}
        _add_table(slide, ["지표", "날짜", "관측값", "기댓값", "Z-score", "방향"],
                   rows, Inches(0.4), Inches(2.9), Inches(12.5), Inches(3.5),
                   highlight_rows=highlight)

        top = max(anomalies, key=lambda a: abs(a.get("z_score", 0)), default=None)
        if top and top.get("llm_interpretation"):
            _add_textbox(slide, f"해석: {top['llm_interpretation']}",
                         Inches(0.4), Inches(6.55), Inches(12.5), Inches(0.55),
                         font_size=10, color=_C["text_mid"], wrap=True)
    else:
        _add_textbox(slide, "이번 주 이상 감지된 지표 없음",
                     Inches(0.4), Inches(3.2), Inches(12.5), Inches(1.0),
                     font_size=18, bold=True, color=_C["positive"], align=PP_ALIGN.CENTER)

        bullets = anomaly_slide.get("bullets", [])
        if bullets:
            y = Inches(4.4)
            for b in bullets[:3]:
                _add_textbox(slide, f"• {b}", Inches(1.5), y, Inches(10.0), Inches(0.42),
                             font_size=12, color=_C["text_dark"])
                y += Inches(0.5)


def _build_slide4_funnel_journey(slide, funnel_metrics: dict, journey_metrics: dict, insight_report: dict) -> None:
    """Layout: [퍼널 테이블 4.2"] [바 차트 3.5"] [journey 경로 5.3"]"""
    _slide_background(slide)
    funnel_slide = (insight_report or {}).get("funnel_slide") or {}
    chart_type = funnel_slide.get("chart_type", "funnel_chart")
    _add_slide_header(slide, "사용자 흐름 분석", funnel_slide.get("headline", ""))

    CONTENT_TOP = Inches(1.3)
    CONTENT_H   = Inches(5.7)   # 슬라이드 높이(7.5) - 헤더(1.15) - 여백

    # ── 왼쪽: 퍼널 테이블 ──────────────────────────────────────────────────────
    steps = (funnel_metrics or {}).get("steps", [])
    if steps:
        _section_label(slide, "퍼널 단계별 전환", Inches(0.3), CONTENT_TOP, width=Inches(4.2))
        funnel_rows = [
            [s.get("event_name", ""),
             _fmt(s.get("user_count", 0)),
             _pct(s.get("conversion_rate", 0)),
             f"{s.get('drop_off_rate', 0):.1f}%"]
            for s in steps
        ]
        _add_table(slide, ["단계", "유저 수", "전환율", "이탈률"],
                   funnel_rows,
                   Inches(0.3), CONTENT_TOP + Inches(0.35),
                   Inches(4.2), CONTENT_H - Inches(0.35),
                   font_size=11)

    # ── 중앙: funnel_chart → 가로 바 차트 ────────────────────────────────────
    if steps:
        step_names = [s.get("event_name", "")[:14] for s in steps]
        user_counts = [float(s.get("user_count", 0) or 0) for s in steps]
        if any(v > 0 for v in user_counts):
            _section_label(slide, "단계별 유저 수", Inches(4.75), CONTENT_TOP, width=Inches(3.5))
            _add_native_chart(
                slide, XL_CHART_TYPE.BAR_CLUSTERED,
                step_names, {"유저 수": user_counts},
                Inches(4.75), CONTENT_TOP + Inches(0.35),
                Inches(3.5), CONTENT_H - Inches(0.35),
            )

    # ── 오른쪽: 전환 경로 + 이탈 패턴 ─────────────────────────────────────────
    converted = (journey_metrics or {}).get("converted_paths", [])[:5]
    if converted:
        _section_label(slide, "주요 전환 경로 Top 5", Inches(8.5), CONTENT_TOP, width=Inches(4.6))
        path_rows = [
            [" → ".join(p.get("path", [])),
             _fmt(p.get("session_count", 0)),
             _pct(p.get("ratio", 0))]
            for p in converted
        ]
        _add_table(slide, ["경로", "세션", "비율"],
                   path_rows,
                   Inches(8.5), CONTENT_TOP + Inches(0.35),
                   Inches(4.6), Inches(3.2),
                   font_size=10)

    journey_summary = (journey_metrics or {}).get("summary", {})
    pre_churn = journey_summary.get("pre_churn_pattern", "")
    if pre_churn:
        _add_rect(slide, Inches(8.5), Inches(5.1), Inches(4.6), Inches(0.75), _C["card"])
        _add_rect(slide, Inches(8.5), Inches(5.1), Inches(0.06), Inches(0.75), _C["negative"])
        _add_textbox(slide, f"이탈 패턴:  {pre_churn}",
                     Inches(8.7), Inches(5.15), Inches(4.2), Inches(0.65),
                     font_size=11, color=_C["negative"])


def _build_slide5_segment(slide, performance_metrics: dict, cohort_metrics: dict, insight_report: dict) -> None:
    _slide_background(slide)
    cohort_slide = (insight_report or {}).get("cohort_slide") or {}
    chart_type   = cohort_slide.get("chart_type", "heatmap")
    _add_slide_header(slide, "고객 세그먼트 분석", cohort_slide.get("headline", ""))

    # ── 상단: 디바이스 / 소스 테이블 (각 절반) ────────────────────────────────
    SEG_TOP = Inches(1.25)
    SEG_H   = Inches(2.5)

    by_device = (performance_metrics or {}).get("by_device_category", [])[:5]
    if by_device:
        _section_label(slide, "디바이스별", Inches(0.4), SEG_TOP, width=Inches(6.0))
        rows = [[d.get("device", ""), _fmt(d.get("session_count", 0)),
                 _pct(d.get("conversion_rate", 0)), _fmt(d.get("revenue", 0), prefix="₩")]
                for d in by_device]
        _add_table(slide, ["디바이스", "세션", "전환율", "매출"],
                   rows, Inches(0.4), SEG_TOP + Inches(0.3), Inches(6.0), SEG_H)

    by_source = (performance_metrics or {}).get("by_traffic_source", [])[:5]
    if by_source:
        _section_label(slide, "트래픽 소스별", Inches(6.8), SEG_TOP, width=Inches(6.2))
        rows = [[s.get("source", "")[:18], _fmt(s.get("session_count", 0)),
                 _pct(s.get("conversion_rate", 0)), _fmt(s.get("revenue", 0), prefix="₩")]
                for s in by_source]
        _add_table(slide, ["소스", "세션", "전환율", "매출"],
                   rows, Inches(6.8), SEG_TOP + Inches(0.3), Inches(6.2), SEG_H)

    # ── 하단: 코호트 히트맵 ────────────────────────────────────────────────────
    COHORT_TOP = Inches(4.05)
    COHORT_H   = Inches(3.0)
    _section_label(slide, "코호트 리텐션 (첫 구매 주차별)", Inches(0.4), COHORT_TOP - Inches(0.3))

    if chart_type == "heatmap":
        _render_cohort_heatmap(slide, cohort_metrics,
                               Inches(0.4), COHORT_TOP, Inches(12.6), COHORT_H)
    else:
        cohort_summary = (cohort_metrics or {}).get("summary", {})
        if cohort_summary:
            cohort_text = (
                f"Week 1 평균 재구매율: {_pct(cohort_summary.get('avg_week1_retention', 0))}  |  "
                f"최고 리텐션 코호트: {cohort_summary.get('best_retention_cohort') or '—'}  |  "
                f"신규 구매자 트렌드: {cohort_summary.get('new_buyer_trend', '—')}"
            )
            _add_rect(slide, Inches(0.4), COHORT_TOP, Inches(12.6), Inches(1.0), _C["card"])
            _add_textbox(slide, cohort_text,
                         Inches(0.6), COHORT_TOP + Inches(0.15), Inches(12.2), Inches(0.7),
                         font_size=12, color=_C["text_dark"], wrap=True, v_anchor="middle")


def _build_slide6_domain(slide, domain: str, performance_metrics: dict, insight_report: dict) -> None:
    _slide_background(slide)
    domain_lower = (domain or "ecommerce").lower().replace("-", "").replace(" ", "")

    if "ecommerce" in domain_lower or "commerce" in domain_lower:
        _add_slide_header(slide, "카테고리별 구매 분석", "e-Commerce 도메인 심화 분석")
        by_category = (performance_metrics or {}).get("by_item_category", [])[:10]
        if by_category:
            rows = [
                [c.get("category", "")[:20],
                 _fmt(c.get("view_count", 0)),
                 _fmt(c.get("add_to_cart_count", 0)),
                 _fmt(c.get("purchase_count", 0)),
                 _fmt(c.get("revenue", 0), prefix="₩"),
                 _pct(c.get("purchase_rate", 0))]
                for c in by_category
            ]
            _add_table(slide, ["카테고리", "조회", "장바구니", "구매", "매출", "구매율"],
                       rows, Inches(0.4), Inches(1.35), Inches(12.5), Inches(5.2))
        else:
            _add_textbox(slide, "(카테고리 데이터 없음)", Inches(0.4), Inches(3.0), Inches(12.5), Inches(1.0),
                         font_size=16, color=_C["text_mid"], align=PP_ALIGN.CENTER)

    elif "fintech" in domain_lower or "finance" in domain_lower:
        _add_slide_header(slide, "금융 상품 분석", "Fintech 도메인 심화 분석")
        _add_textbox(slide, "Fintech 도메인: 금융 상품별 전환 및 리텐션 분석 (구현 예정)",
                     Inches(0.4), Inches(2.5), Inches(12.5), Inches(2.0),
                     font_size=16, color=_C["text_mid"], align=PP_ALIGN.CENTER)

    elif "media" in domain_lower or "content" in domain_lower:
        _add_slide_header(slide, "콘텐츠 소비 분석", "Media 도메인 심화 분석")
        _add_textbox(slide, "Media 도메인: 콘텐츠 유형별 참여율 및 완독률 분석 (구현 예정)",
                     Inches(0.4), Inches(2.5), Inches(12.5), Inches(2.0),
                     font_size=16, color=_C["text_mid"], align=PP_ALIGN.CENTER)

    else:
        _add_slide_header(slide, f"{domain} 도메인 분석", "도메인 특화 분석")
        _add_textbox(slide,
                     f"'{domain}' 도메인의 특화 분석 슬라이드입니다.\n"
                     "도메인 컨텍스트에 맞게 이 슬라이드를 커스터마이징하세요.",
                     Inches(0.4), Inches(2.5), Inches(12.5), Inches(2.0),
                     font_size=15, color=_C["text_mid"], wrap=True)


def _build_slide7_prediction(slide, prediction_metrics: dict, insight_report: dict) -> None:
    _slide_background(slide)
    pred_slide = (insight_report or {}).get("prediction_slide") or {}
    chart_type = pred_slide.get("chart_type", "line_chart")
    _add_slide_header(slide, "예측 및 시사점", pred_slide.get("headline", ""))

    predictions = (prediction_metrics or {}).get("predictions", [])
    pred_summary = (prediction_metrics or {}).get("summary", {})
    overall_trend = pred_summary.get("overall_trend", "stable")
    dq_warning = pred_summary.get("data_quality_warning")

    trend_color = {"increasing": _C["positive"], "decreasing": _C["negative"]}.get(
        overall_trend, _C["neutral"])
    _add_rect(slide, Inches(0.4), Inches(1.25), Inches(4.5), Inches(0.6), _C["card"])
    _add_rect(slide, Inches(0.4), Inches(1.25), Inches(0.06), Inches(0.6), trend_color)
    _add_textbox(slide, f"전반적 추세: {overall_trend}",
                 Inches(0.6), Inches(1.3), Inches(4.2), Inches(0.5),
                 font_size=14, bold=True, color=trend_color)

    if dq_warning:
        _add_textbox(slide, f"⚠ {dq_warning}",
                     Inches(5.2), Inches(1.35), Inches(7.8), Inches(0.45),
                     font_size=11, color=_C["neutral"])

    if predictions:
        rows = []
        for p in predictions:
            ci = p.get("confidence_interval", {})
            skipped = p.get("skipped", False)
            rows.append([
                p.get("target", ""),
                _fmt(p.get("predicted_value", 0), decimals=2) if not skipped else "데이터 부족",
                (f"{_fmt(ci.get('lower', 0), decimals=2)} ~ {_fmt(ci.get('upper', 0), decimals=2)}"
                 if not skipped else "—"),
                p.get("trend_direction", ""),
                (p.get("llm_comment") or "—")[:60],
            ])

        has_chart_data = any(not p.get("skipped") for p in predictions)

        if has_chart_data and chart_type == "line_chart":
            # Left: prediction table (narrower); Right: line chart
            _add_table(slide, ["예측 지표", "예측값", "신뢰구간", "추세", "코멘트"],
                       rows, Inches(0.4), Inches(2.0), Inches(7.0), Inches(3.8))
            _render_chart_for_slide(
                slide, chart_type, prediction_metrics,
                Inches(7.7), Inches(2.0), Inches(5.3), Inches(3.8),
            )
        else:
            _add_table(slide, ["예측 지표", "예측값", "신뢰구간", "추세", "코멘트"],
                       rows, Inches(0.4), Inches(2.0), Inches(12.5), Inches(3.8))

    bullets = pred_slide.get("bullets", [])
    if bullets:
        y = Inches(6.0)
        for b in bullets[:2]:
            _add_textbox(slide, f"• {b}", Inches(0.4), y, Inches(12.5), Inches(0.38),
                         font_size=10, color=_C["text_mid"])
            y += Inches(0.4)


def _build_slide8_recommendations(slide, insight_report: dict) -> None:
    _slide_background(slide)
    _add_slide_header(slide, "권장 액션", "컨텍스트 기반 우선순위 권장사항")

    recommendations = (insight_report or {}).get("recommendations", [])
    cross_findings = (insight_report or {}).get("cross_analysis_findings", [])

    priority_colors = [_C["negative"], _C["neutral"], _C["positive"]]
    priority_labels = ["P1", "P2", "P3"]

    if recommendations:
        _section_label(slide, "우선순위 권장 액션", Inches(0.4), Inches(1.25))
        y = Inches(1.65)
        for i, rec in enumerate(recommendations[:5]):
            color = priority_colors[min(i, 2)]
            label = priority_labels[min(i, 2)]
            _add_rect(slide, Inches(0.4), y, Inches(0.55), Inches(0.5), color)
            _add_textbox(slide, label, Inches(0.4), y, Inches(0.55), Inches(0.5),
                         font_size=11, bold=True, color=_C["text_light"],
                         align=PP_ALIGN.CENTER, v_anchor="middle")
            _add_rect(slide, Inches(1.05), y, Inches(12.0), Inches(0.5), _C["card"])
            _add_textbox(slide, rec, Inches(1.05), y, Inches(12.0), Inches(0.5),
                         font_size=12, color=_C["text_dark"], wrap=True, v_anchor="middle")
            y += Inches(0.65)

    if cross_findings:
        _section_label(slide, "교차 분석 인사이트", Inches(0.4), Inches(5.1))
        y = Inches(5.45)
        for finding in cross_findings[:3]:
            _add_rect(slide, Inches(0.4), y, Inches(12.5), Inches(0.42), _C["card"])
            _add_rect(slide, Inches(0.4), y, Inches(0.06), Inches(0.42), _C["accent"])
            _add_textbox(slide, f"  {finding}", Inches(0.55), y, Inches(12.2), Inches(0.42),
                         font_size=11, color=_C["text_dark"], v_anchor="middle")
            y += Inches(0.52)


# ──────────────────────────────────────────────────────────────────────────────
# Main builder
# ──────────────────────────────────────────────────────────────────────────────

def _build_presentation(state: dict) -> Presentation:
    insight_report      = state.get("insight_report") or {}
    performance_metrics = state.get("performance_metrics") or {}
    funnel_metrics      = state.get("funnel_metrics") or {}
    cohort_metrics      = state.get("cohort_metrics") or {}
    journey_metrics     = state.get("journey_metrics") or {}
    anomaly_metrics     = state.get("anomaly_metrics") or {}
    prediction_metrics  = state.get("prediction_metrics") or {}
    domain_context      = state.get("domain_context") or {}
    domain              = domain_context.get("domain", "ecommerce")

    prs = Presentation()
    prs.slide_width  = _SLIDE_W
    prs.slide_height = _SLIDE_H
    blank = prs.slide_layouts[6]

    builders = [
        lambda s: _build_slide1_executive_summary(s, insight_report, performance_metrics),
        lambda s: _build_slide2_performance(s, performance_metrics, insight_report),
        lambda s: _build_slide3_anomaly(s, anomaly_metrics, insight_report),
        lambda s: _build_slide4_funnel_journey(s, funnel_metrics, journey_metrics, insight_report),
        lambda s: _build_slide5_segment(s, performance_metrics, cohort_metrics, insight_report),
        lambda s: _build_slide6_domain(s, domain, performance_metrics, insight_report),
        lambda s: _build_slide7_prediction(s, prediction_metrics, insight_report),
        lambda s: _build_slide8_recommendations(s, insight_report),
    ]
    for fn in builders:
        slide = prs.slides.add_slide(blank)
        fn(slide)
        _add_logo(slide)  # always last → on top of everything

    return prs


# ──────────────────────────────────────────────────────────────────────────────
# Agent entry point
# ──────────────────────────────────────────────────────────────────────────────

async def ppt_agent(state: dict) -> dict:
    """LangGraph node: generate an 8-slide PowerPoint report."""
    if not state.get("insight_report"):
        raise ValueError("ppt_agent: 'insight_report' must be present in pipeline state")

    week_start = state.get("week_start", "unknown")
    week_end   = state.get("week_end", "unknown")
    domain     = (state.get("domain_context") or {}).get("domain", "report")

    prs = _build_presentation(state)

    os.makedirs(_OUTPUT_DIR, exist_ok=True)
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    filename  = f"{domain}_{week_start}_{week_end}_{timestamp}.pptx"
    filepath  = os.path.join(_OUTPUT_DIR, filename)
    prs.save(filepath)

    logger.info("ppt_agent: saved report → %s", filepath)
    return {"ppt_url": filepath}
